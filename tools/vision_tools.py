##########################################################################################
#
# Module: tools/vision_tools.py
#
# Description: Vision tools for agent use.
#              Provides image analysis and document extraction capabilities.
#
# Author: Cornelis Networks
#
##########################################################################################

import base64
import logging
import os
import sys
from pathlib import Path
from typing import Any, Dict, List, Optional

from tools.base import BaseTool, ToolResult, tool

# Logging config - follows jira_utils.py pattern
log = logging.getLogger(os.path.basename(sys.argv[0]))

# Optional imports for document processing
try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False
    log.debug('PIL not available - image processing limited')

try:
    import openpyxl
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False
    log.debug('openpyxl not available - Excel processing disabled')

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    log.debug('python-pptx not available - PowerPoint processing disabled')


# ****************************************************************************************
# Tool Functions
# ****************************************************************************************

@tool(
    description='Analyze an image using vision LLM to extract information'
)
def analyze_image(
    image_path: str,
    prompt: str = 'Describe this image in detail.',
    extract_text: bool = True
) -> ToolResult:
    '''
    Analyze an image using a vision-capable LLM.
    
    Input:
        image_path: Path to the image file.
        prompt: The analysis prompt/question about the image.
        extract_text: Whether to also extract any visible text.
    
    Output:
        ToolResult with analysis results including:
        - description: LLM's description of the image
        - extracted_text: Any text found in the image (if extract_text=True)
        - metadata: Image metadata (size, format, etc.)
    '''
    log.debug(f'analyze_image(image_path={image_path})')
    
    try:
        if not os.path.exists(image_path):
            return ToolResult.failure(f'Image file not found: {image_path}')
        
        # Get image metadata
        metadata = _get_image_metadata(image_path)
        
        # Convert image to base64 for LLM
        image_data = _image_to_base64(image_path)
        
        # Import LLM client
        from llm.config import get_vision_client
        
        try:
            llm = get_vision_client()
        except Exception as e:
            log.warning(f'Vision LLM not available: {e}')
            # Return just metadata if no vision LLM
            return ToolResult.success({
                'description': 'Vision LLM not available for analysis',
                'metadata': metadata,
                'image_data': image_data[:100] + '...'  # Truncated for display
            })
        
        # Build the analysis prompt
        full_prompt = prompt
        if extract_text:
            full_prompt += '\n\nAlso extract and list any text visible in the image.'
        
        # Call vision LLM
        from llm.base import Message
        messages = [Message.user(full_prompt)]
        
        response = llm.chat_with_vision(
            messages=messages,
            images=[f'data:image/{metadata.get("format", "png")};base64,{image_data}'],
            max_tokens=2000
        )
        
        result = {
            'description': response.content,
            'metadata': metadata,
            'model_used': response.model,
            'tokens_used': response.total_tokens
        }
        
        return ToolResult.success(result)
        
    except Exception as e:
        log.error(f'Failed to analyze image: {e}')
        return ToolResult.failure(f'Image analysis failed: {e}')


@tool(
    description='Extract roadmap information from a PowerPoint presentation'
)
def extract_roadmap_from_ppt(
    file_path: str,
    slide_numbers: Optional[List[int]] = None
) -> ToolResult:
    '''
    Extract roadmap information from a PowerPoint file.
    
    Parses the PowerPoint to extract text, shapes, and structure,
    then uses vision LLM to interpret roadmap content.
    
    Input:
        file_path: Path to the .pptx file.
        slide_numbers: Optional list of specific slide numbers to process (1-indexed).
    
    Output:
        ToolResult with extracted roadmap data:
        - slides: List of slide content
        - releases: Detected release information
        - timeline: Detected timeline/milestones
        - features: Detected features/items
    '''
    log.debug(f'extract_roadmap_from_ppt(file_path={file_path})')
    
    if not PPTX_AVAILABLE:
        return ToolResult.failure('python-pptx package not installed. Run: pip install python-pptx')
    
    try:
        if not os.path.exists(file_path):
            return ToolResult.failure(f'File not found: {file_path}')
        
        prs = Presentation(file_path)
        
        slides_data = []
        all_text = []
        
        for i, slide in enumerate(prs.slides, 1):
            # Skip if not in requested slides
            if slide_numbers and i not in slide_numbers:
                continue
            
            slide_data = {
                'number': i,
                'title': '',
                'text_content': [],
                'shapes': []
            }
            
            for shape in slide.shapes:
                # Extract title
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        slide_data['text_content'].append(text)
                        all_text.append(text)
                        
                        # Check if this is the title
                        if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                            if shape.placeholder_format.type == 1:  # Title placeholder
                                slide_data['title'] = text
                
                # Extract shape info
                if hasattr(shape, 'shape_type'):
                    slide_data['shapes'].append({
                        'type': str(shape.shape_type),
                        'has_text': shape.has_text_frame
                    })
            
            slides_data.append(slide_data)
        
        # Try to parse roadmap structure from text
        roadmap_data = _parse_roadmap_text(all_text)
        
        result = {
            'slides': slides_data,
            'slide_count': len(slides_data),
            'total_slides': len(prs.slides),
            'releases': roadmap_data.get('releases', []),
            'timeline': roadmap_data.get('timeline', []),
            'features': roadmap_data.get('features', []),
            'raw_text': all_text
        }
        
        return ToolResult.success(result)
        
    except Exception as e:
        log.error(f'Failed to extract from PowerPoint: {e}')
        return ToolResult.failure(f'PowerPoint extraction failed: {e}')


@tool(
    description='Extract roadmap information from an Excel spreadsheet'
)
def extract_roadmap_from_excel(
    file_path: str,
    sheet_name: Optional[str] = None
) -> ToolResult:
    '''
    Extract roadmap information from an Excel file.
    
    Parses the Excel file to extract structured data that may
    represent a roadmap or release plan.
    
    Input:
        file_path: Path to the .xlsx file.
        sheet_name: Optional specific sheet to process.
    
    Output:
        ToolResult with extracted data:
        - sheets: List of sheet data
        - releases: Detected release information
        - timeline: Detected timeline/milestones
        - features: Detected features/items
    '''
    log.debug(f'extract_roadmap_from_excel(file_path={file_path})')
    
    if not OPENPYXL_AVAILABLE:
        return ToolResult.failure('openpyxl package not installed. Run: pip install openpyxl')
    
    try:
        if not os.path.exists(file_path):
            return ToolResult.failure(f'File not found: {file_path}')
        
        wb = openpyxl.load_workbook(file_path, data_only=True)
        
        sheets_data = []
        all_data = []
        
        for ws_name in wb.sheetnames:
            if sheet_name and ws_name != sheet_name:
                continue
            
            ws = wb[ws_name]
            
            sheet_data = {
                'name': ws_name,
                'rows': [],
                'headers': []
            }
            
            # Get headers from first row
            headers = []
            for cell in ws[1]:
                headers.append(str(cell.value) if cell.value else '')
            sheet_data['headers'] = headers
            
            # Get data rows
            for row in ws.iter_rows(min_row=2, values_only=True):
                row_data = {}
                for i, value in enumerate(row):
                    if i < len(headers) and headers[i]:
                        row_data[headers[i]] = value
                    else:
                        row_data[f'col_{i}'] = value
                
                if any(v is not None for v in row_data.values()):
                    sheet_data['rows'].append(row_data)
                    all_data.append(row_data)
            
            sheets_data.append(sheet_data)
        
        # Try to identify roadmap columns
        roadmap_data = _parse_roadmap_excel(sheets_data)
        
        result = {
            'sheets': sheets_data,
            'sheet_count': len(sheets_data),
            'releases': roadmap_data.get('releases', []),
            'timeline': roadmap_data.get('timeline', []),
            'features': roadmap_data.get('features', [])
        }
        
        return ToolResult.success(result)
        
    except Exception as e:
        log.error(f'Failed to extract from Excel: {e}')
        return ToolResult.failure(f'Excel extraction failed: {e}')


@tool(
    description='Extract text from an image using OCR or vision LLM'
)
def extract_text_from_image(image_path: str) -> ToolResult:
    '''
    Extract text content from an image.
    
    Uses vision LLM to perform OCR-like text extraction.
    
    Input:
        image_path: Path to the image file.
    
    Output:
        ToolResult with extracted text.
    '''
    log.debug(f'extract_text_from_image(image_path={image_path})')
    
    # Use analyze_image with a text extraction prompt
    return analyze_image(
        image_path=image_path,
        prompt='Extract all text visible in this image. Return the text exactly as it appears, preserving layout where possible.',
        extract_text=True
    )


# ****************************************************************************************
# Helper Functions
# ****************************************************************************************

def _get_image_metadata(image_path: str) -> Dict[str, Any]:
    '''Get metadata about an image file.'''
    metadata = {
        'path': image_path,
        'filename': os.path.basename(image_path),
        'size_bytes': os.path.getsize(image_path),
        'format': Path(image_path).suffix.lower().lstrip('.')
    }
    
    if PIL_AVAILABLE:
        try:
            with Image.open(image_path) as img:
                metadata['width'] = img.width
                metadata['height'] = img.height
                metadata['mode'] = img.mode
                metadata['format'] = img.format or metadata['format']
        except Exception as e:
            log.debug(f'Could not read image with PIL: {e}')
    
    return metadata


def _image_to_base64(image_path: str) -> str:
    '''Convert an image file to base64 string.'''
    with open(image_path, 'rb') as f:
        return base64.b64encode(f.read()).decode('utf-8')


def _parse_roadmap_text(text_items: List[str]) -> Dict[str, List]:
    '''
    Parse roadmap information from text items.
    
    Looks for patterns indicating releases, dates, and features.
    '''
    releases = []
    timeline = []
    features = []
    
    import re
    
    # Patterns for release versions
    version_pattern = re.compile(r'\b(\d+\.\d+(?:\.\d+)?)\b')
    # Patterns for dates
    date_pattern = re.compile(r'\b(Q[1-4]\s*\d{4}|\d{4}\s*Q[1-4]|(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]*\s*\d{4})\b', re.IGNORECASE)
    
    for text in text_items:
        # Look for version numbers
        versions = version_pattern.findall(text)
        for v in versions:
            if v not in [r['version'] for r in releases]:
                releases.append({'version': v, 'context': text[:100]})
        
        # Look for dates
        dates = date_pattern.findall(text)
        for d in dates:
            timeline.append({'date': d, 'context': text[:100]})
        
        # Treat other text as potential features
        if len(text) > 10 and not versions and not dates:
            features.append({'text': text[:200]})
    
    return {
        'releases': releases,
        'timeline': timeline,
        'features': features[:20]  # Limit features
    }


def _parse_roadmap_excel(sheets_data: List[Dict]) -> Dict[str, List]:
    '''
    Parse roadmap information from Excel data.
    
    Looks for columns that might contain release, date, or feature info.
    '''
    releases = []
    timeline = []
    features = []
    
    # Common column names for roadmap data
    release_columns = ['release', 'version', 'milestone', 'target']
    date_columns = ['date', 'due', 'target', 'eta', 'quarter', 'q1', 'q2', 'q3', 'q4']
    feature_columns = ['feature', 'item', 'description', 'name', 'title', 'task', 'story']
    
    for sheet in sheets_data:
        headers_lower = [h.lower() for h in sheet.get('headers', [])]
        
        # Find relevant columns
        release_col = None
        date_col = None
        feature_col = None
        
        for i, h in enumerate(headers_lower):
            if any(rc in h for rc in release_columns):
                release_col = sheet['headers'][i]
            if any(dc in h for dc in date_columns):
                date_col = sheet['headers'][i]
            if any(fc in h for fc in feature_columns):
                feature_col = sheet['headers'][i]
        
        # Extract data from rows
        for row in sheet.get('rows', []):
            if release_col and row.get(release_col):
                releases.append({
                    'version': str(row[release_col]),
                    'sheet': sheet['name']
                })
            
            if date_col and row.get(date_col):
                timeline.append({
                    'date': str(row[date_col]),
                    'sheet': sheet['name']
                })
            
            if feature_col and row.get(feature_col):
                features.append({
                    'text': str(row[feature_col]),
                    'sheet': sheet['name']
                })
    
    return {
        'releases': releases,
        'timeline': timeline,
        'features': features
    }


# ****************************************************************************************
# Tool Collection Class
# ****************************************************************************************

class VisionTools(BaseTool):
    '''
    Collection of vision tools for agent use.
    '''
    
    @tool(description='Analyze an image using vision LLM')
    def analyze_image(
        self,
        image_path: str,
        prompt: str = 'Describe this image in detail.'
    ) -> ToolResult:
        return analyze_image(image_path, prompt)
    
    @tool(description='Extract roadmap from PowerPoint')
    def extract_roadmap_from_ppt(
        self,
        file_path: str,
        slide_numbers: Optional[List[int]] = None
    ) -> ToolResult:
        return extract_roadmap_from_ppt(file_path, slide_numbers)
    
    @tool(description='Extract roadmap from Excel')
    def extract_roadmap_from_excel(
        self,
        file_path: str,
        sheet_name: Optional[str] = None
    ) -> ToolResult:
        return extract_roadmap_from_excel(file_path, sheet_name)
    
    @tool(description='Extract text from an image')
    def extract_text_from_image(self, image_path: str) -> ToolResult:
        return extract_text_from_image(image_path)
